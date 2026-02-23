from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
import copy
from lxml import etree

# ── Color constants ──────────────────────────────────────────────
BG     = RGBColor(15,  17,  22)
CARD   = RGBColor(25,  28,  36)
LIME   = RGBColor(200, 241, 53)
BLUE   = RGBColor(91,  143, 255)
WHITE  = RGBColor(255, 255, 255)
GRAY   = RGBColor(136, 136, 136)
ORANGE = RGBColor(255, 107, 53)
GREEN2 = RGBColor(52,  211, 153)
PURPLE = RGBColor(167, 139, 250)
BLACK  = RGBColor(0,   0,   0)

# ── Presentation setup ───────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

# ── Helper functions ─────────────────────────────────────────────

def set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, text="", font_size=14,
             font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border by default

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = "Calibri"
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=14, font_color=WHITE,
                bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Calibri"
    return txBox


def add_multiline_textbox(slide, x, y, w, h, lines, font_size=12,
                           font_color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                           line_spacing_pt=None):
    """Add a textbox with multiple lines (list of strings)."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = "Calibri"
        if line_spacing_pt:
            from pptx.util import Pt as Pt2
            p.line_spacing = Pt2(line_spacing_pt)
    return txBox


def add_title(slide, text):
    return add_textbox(slide, 0.5, 0.3, 12.3, 0.7, text,
                       font_size=32, font_color=LIME, bold=True,
                       align=PP_ALIGN.LEFT)


def add_card(slide, x, y, w, h, title, body_lines, title_color=LIME,
             body_font_size=11):
    # Card background
    add_rect(slide, x, y, w, h, CARD)
    # Title
    add_textbox(slide, x + 0.1, y + 0.1, w - 0.2, 0.35, title,
                font_size=13, font_color=title_color, bold=True,
                align=PP_ALIGN.LEFT)
    # Body lines
    add_multiline_textbox(slide, x + 0.1, y + 0.5, w - 0.2, h - 0.6,
                          body_lines, font_size=body_font_size,
                          font_color=WHITE, bold=False, align=PP_ALIGN.LEFT)
    return


def add_left_border(slide, x, y, w, h, border_color, thickness=0.06):
    """Draw a thin colored rectangle as a left border accent."""
    add_rect(slide, x, y, thickness, h, border_color)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide1, BG)

# Bottom bar
add_rect(slide1, 0, 6.5, 13.33, 1.0, LIME)

# "28th Project" badge
add_rect(slide1, 0.5, 0.2, 2.0, 0.4, LIME, "28th Project",
         font_size=12, font_color=BLACK, bold=True)

# Main title
add_textbox(slide1, 1, 1.5, 11.33, 2.0, "PoseCoach",
            font_size=80, font_color=LIME, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1, 1, 3.3, 11.33, 0.7,
            "AI 기반 운동 자세 코칭 시스템",
            font_size=28, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)

# Tags
add_textbox(slide1, 1, 4.0, 11.33, 0.5,
            "YOLO Pose  ·  DTW  ·  Gemini AI  ·  FastAPI  ·  React",
            font_size=16, font_color=GRAY, bold=False, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide2, BG)
add_title(slide2, "왜 자세 코칭이 필요한가?")

card_y  = 1.5
card_h  = 2.5
card_w  = 3.5

cards = [
    (0.5,  "❌ 자세 확인 불가",
     ["혼자 운동 시", "자세 피드백 없음", "→ 잘못된 습관 형성"]),
    (4.4,  "❌ 부상 위험",
     ["잘못된 자세 반복", "→ 관절 과부하", "→ 만성 부상"]),
    (8.3,  "❌ 동기 부족",
     ["성장이 보이지 않아", "운동 지속 어려움", "→ 포기"]),
]

for cx, ctitle, cbody in cards:
    # border
    add_rect(slide2, cx - 0.05, card_y - 0.05, card_w + 0.1, card_h + 0.1, ORANGE)
    add_rect(slide2, cx, card_y, card_w, card_h, CARD)
    add_textbox(slide2, cx + 0.15, card_y + 0.15, card_w - 0.3, 0.45,
                ctitle, font_size=14, font_color=ORANGE, bold=True)
    add_multiline_textbox(slide2, cx + 0.15, card_y + 0.7, card_w - 0.3,
                          card_h - 0.9, cbody, font_size=13,
                          font_color=WHITE)

# Solution banner
add_rect(slide2, 0.5, 4.3, 12.3, 0.8, LIME,
         "→ PoseCoach: 영상 하나로 AI가 자세 분석 · 점수 · 피드백 제공",
         font_size=18, font_color=BLACK, bold=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide3, BG)
add_title(slide3, "시스템 전체 구조")

layers = [
    ((30, 35, 50),  "👤 사용자",                          WHITE,  "영상 업로드 / 결과 확인",                              GRAY),
    ((20, 40, 60),  "🌐 Frontend (React + TypeScript)",    BLUE,   "UploadVideo · Result · MyPage · Home",                 GRAY),
    ((20, 50, 40),  "⚙️ Backend (FastAPI)",                GREEN2, "분석 API · 인증 API · DB · Gemini 피드백 API",          GRAY),
    ((50, 30, 20),  "🤖 AI 분석 엔진 (Python)",            ORANGE, "YOLO Pose → Phase 감지 → 횟수 카운팅 → 자세 평가 → DTW", GRAY),
    ((40, 20, 50),  "🗄️ Database (SQLite)",               PURPLE, "유저 · 운동기록 · Phase점수 · 오류프레임 저장",           GRAY),
]

y_start = 1.3
box_h   = 0.82
gap     = 0.04

for i, (rgb, left_text, left_color, right_text, right_color) in enumerate(layers):
    y = y_start + i * (box_h + gap + 0.18)
    fill_c = RGBColor(*rgb)
    # Full-width background
    add_rect(slide3, 0.65, y, 12.0, box_h, fill_c)
    # Left label (20% width)
    add_textbox(slide3, 0.75, y + 0.1, 3.2, box_h - 0.2, left_text,
                font_size=13, font_color=left_color, bold=True)
    # Right detail (80% width)
    add_textbox(slide3, 4.1, y + 0.1, 8.4, box_h - 0.2, right_text,
                font_size=11, font_color=right_color, bold=False)
    # Arrow between layers
    if i < len(layers) - 1:
        add_textbox(slide3, 6.3, y + box_h, 0.8, 0.22, "↓",
                    font_size=12, font_color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide4, BG)
add_title(slide4, "기술 스택")

tech_cols = [
    (0.3,  "AI / ML",      LIME,   [
        "• YOLO v8 Pose", "• fastdtw", "• Cohen's d",
        "• Gemini API",   "• scikit-learn", "• OpenCV"]),
    (3.5,  "Backend",      BLUE,   [
        "• Python 3.11", "• FastAPI", "• SQLite",
        "• bcrypt",      "• pydantic", "• uvicorn"]),
    (6.7,  "Frontend",     GREEN2, [
        "• React 18",    "• TypeScript", "• Vite",
        "• Tailwind CSS","• React Router","• shadcn/ui"]),
    (9.9,  "Infra / Tools",ORANGE, [
        "• Git / GitHub","• Streamlit", "• Plotly",
        "• python-pptx", "• reportlab", "• SQLite WAL"]),
]

for cx, ctitle, ccolor, clines in tech_cols:
    add_card(slide4, cx, 1.3, 2.8, 4.5, ctitle, clines,
             title_color=ccolor, body_font_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Video Processing Pipeline
# ════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide5, BG)
add_title(slide5, "영상 처리 파이프라인")

steps = [
    # (x, y, header_color, title, body_lines)
    (0.3,  1.5, LIME,   "① 영상 업로드",
     ["MP4/MOV/AVI/WEBM", "최대 1920×1080 해상도"]),
    (4.5,  1.5, BLUE,   "② 프레임 추출",
     ["설정 FPS(1~30)로", "이미지 추출", "(default: 10fps)"]),
    (8.7,  1.5, GREEN2, "③ 활성 구간 필터",
     ["ML + Rule-based", "하이브리드 필터", "운동 구간만 선택"]),
    (0.3,  3.3, ORANGE, "④ YOLO 포즈 추정",
     ["17 keypoints", "+ 가상 3개", "= 총 20개"]),
    (4.5,  3.3, PURPLE, "⑤ 키포인트 스무딩",
     ["window=3", "이동평균 필터", "떨림 제거"]),
    (8.7,  3.3, LIME,   "⑥ 정규화 좌표",
     ["[0,1] 범위 변환", "카메라 해상도", "무관"]),
]

for sx, sy, hcolor, stitle, sbody in steps:
    add_rect(slide5, sx, sy, 3.8, 1.5, CARD)
    add_rect(slide5, sx, sy, 3.8, 0.38, hcolor)
    add_textbox(slide5, sx + 0.1, sy + 0.04, 3.6, 0.32,
                stitle, font_size=12, font_color=BLACK, bold=True)
    add_multiline_textbox(slide5, sx + 0.1, sy + 0.44, 3.6, 0.95,
                          sbody, font_size=11, font_color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Activity Filter
# ════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide6, BG)
add_title(slide6, "활성 구간 필터링 (운동 구간 감지)")

# Left: problem & solution
add_rect(slide6, 0.4, 1.2, 6.0, 0.7, CARD,
         "문제: 전체 영상 분석 시 준비/휴식 구간 포함 → 낮은 점수, 왜곡된 분석 결과",
         font_size=11, font_color=ORANGE, bold=True, align=PP_ALIGN.LEFT)

add_textbox(slide6, 0.4, 2.05, 6.0, 0.35,
            "해결 방법 — 하이브리드 필터",
            font_size=13, font_color=LIME, bold=True)

filter_steps = [
    (BLUE,   "① ML 모델",
     ["(activity_filter.pkl)", "RF 분류기로", "활성 확률 예측"]),
    (GREEN2, "② Rule-based",
     ["모션 감지", "프레임 간 차이", "임계값 필터링"]),
    (LIME,   "③ 하이브리드 앙상블",
     ["ML OR (rule AND prob)", "→ 최종 활성 구간 선택"]),
]

for i, (fc, ft, fb) in enumerate(filter_steps):
    bx = 0.4 + i * 2.05
    add_rect(slide6, bx, 2.5, 1.9, 2.2, CARD)
    add_rect(slide6, bx, 2.5, 1.9, 0.38,
             fc if fc != LIME else LIME)
    txt_color = BLACK if fc == LIME else BLACK
    add_textbox(slide6, bx + 0.05, 2.54, 1.8, 0.32,
                ft, font_size=10, font_color=txt_color if fc == LIME else WHITE, bold=True)
    add_multiline_textbox(slide6, bx + 0.05, 2.94, 1.8, 1.7,
                          fb, font_size=10, font_color=WHITE)

# Right: stats
add_textbox(slide6, 6.8, 1.2, 6.0, 0.4,
            "필터링 효과", font_size=14, font_color=LIME, bold=True)

stats_lines = [
    "• 전체 프레임: 100%",
    "• 필터 후: ~35-60%",
    "• 분석 속도: 2-3배 향상",
    "• 점수 왜곡 제거",
]
add_multiline_textbox(slide6, 6.8, 1.7, 6.0, 1.4,
                      stats_lines, font_size=12, font_color=WHITE)

add_textbox(slide6, 6.8, 3.2, 6.0, 0.4,
            "방법 비교", font_size=13, font_color=BLUE, bold=True)

method_lines = [
    "ML 방법: O(N) 추론, 높은 정확도",
    "Rule 방법: 빠른 fallback",
    "Gap fill: 짧은 정지 구간 보완",
]
add_multiline_textbox(slide6, 6.8, 3.7, 6.0, 1.2,
                      method_lines, font_size=12, font_color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Phase Detection & Counter
# ════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide7, BG)
add_title(slide7, "Phase 감지 & 운동 횟수 카운팅")

# Phase timeline bar
phases = [
    ("ready",      GRAY),
    ("top",        LIME),
    ("descending", BLUE),
    ("bottom",     ORANGE),
    ("ascending",  GREEN2),
]
bar_x = 0.5
bar_y = 1.2
bar_h = 0.55
seg_w = 12.3 / len(phases)
for i, (pname, pcolor) in enumerate(phases):
    add_rect(slide7, bar_x + i * seg_w, bar_y, seg_w - 0.02, bar_h,
             pcolor, pname, font_size=12,
             font_color=BLACK if pcolor == LIME else WHITE, bold=True)

# Left: Phase detection algorithm
add_card(slide7, 0.4, 2.1, 5.8, 4.2,
         "Phase 감지 알고리즘", [
             "• 팔꿈치 각도 + 각속도 기반 감지",
             "• FPS 적응형 임계값 (fps 스케일링)",
             "• Hysteresis: 노이즈 방지",
             "• 속도 스무딩: jitter 제거",
             "",
             "푸시업 기준:",
             "  top: 팔꿈치 > 150°",
             "  bottom: 팔꿈치 < 110°",
             "",
             "풀업 기준:",
             "  bottom: 팔꿈치 > 140°",
             "  top: 팔꿈치 < 50°",
         ], title_color=LIME, body_font_size=11)

# Right: Counter logic
add_card(slide7, 6.7, 2.1, 5.8, 4.2,
         "횟수 카운팅 로직", [
             "• 활성화 조건 감지",
             "  푸시업: 손목 < 어깨 높이",
             "  풀업: 손목 > 어깨 높이",
             "",
             "• Rep 완료 조건:",
             "  필요 Phase 순서 완주",
             "  → count += 1",
             "",
             "• 비활성 타임아웃:",
             "  FPS 기반 자동 리셋",
             "  (2초 이상 정지 감지)",
         ], title_color=BLUE, body_font_size=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Scoring System
# ════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide8, BG)
add_title(slide8, "자세 평가 시스템 (점수 로직)")

# Badge
add_rect(slide8, 10.5, 0.25, 2.3, 0.4, LIME,
         "👤 담당 파트", font_size=12, font_color=BLACK, bold=True)

score_cards = [
    (0.3,   "① Cohen's d 가중치", LIME, [
        "AI Hub 데이터 224개 분석",
        "효과크기 기반 중요도 산출",
        "",
        "어깨외전: 0.32",
        "손 위치:  0.28",
        "고개 숙임: 0.24",
        "팔꿈치:   0.12",
        "등 직선:  0.04",
    ]),
    (4.55,  "② Soft Scoring",     BLUE, [
        "이진(0/1) → 연속(0~1)",
        "",
        "임계값 근방 선형 보간",
        "미세한 자세 차이 반영",
        "",
        "예: 팔꿈치 각도",
        "159° → 0.95점",
        "140° → 0.25점",
        "100° → 0.00점",
    ]),
    (8.8,   "③ DTW 유사도",        GREEN2, [
        "레퍼런스 영상과 비교",
        "7개 각도 피처 사용",
        "fastdtw O(N) 알고리즘",
        "",
        "Combined Score =",
        "avg × 0.7",
        "+ dtw × 0.3",
        "",
        "→ S/A/B/C 등급",
    ]),
]

for cx, ctitle, ccolor, clines in score_cards:
    add_card(slide8, cx, 1.4, 3.8, 4.6, ctitle, clines,
             title_color=ccolor, body_font_size=11)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Backend API
# ════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide9, BG)
add_title(slide9, "백엔드 API (FastAPI)")

# Left: endpoints
add_textbox(slide9, 0.4, 1.2, 5.5, 0.4,
            "API 엔드포인트", font_size=14, font_color=LIME, bold=True)

endpoint_lines = [
    "POST  /auth/register  →  회원가입",
    "POST  /auth/login     →  JWT 로그인",
    "POST  /analyze        →  영상 분석",
    "GET   /history        →  운동 기록",
    "POST  /report/pdf     →  PDF 리포트",
    "POST  /ai-feedback    →  Gemini 피드백",
]
add_rect(slide9, 0.4, 1.7, 5.8, 4.5, CARD)
add_multiline_textbox(slide9, 0.55, 1.85, 5.5, 4.2,
                      endpoint_lines, font_size=11.5, font_color=WHITE)

# Right: flow
add_textbox(slide9, 6.5, 1.2, 6.4, 0.4,
            "분석 API 흐름", font_size=14, font_color=BLUE, bold=True)

flow_items = [
    ("영상 파일 수신 + 메타데이터",  CARD,  WHITE),
    ("↓",                            BG,    GRAY),
    ("run_video_analysis() 실행",     CARD,  WHITE),
    ("↓",                            BG,    GRAY),
    ("frame_scores / error_frames 반환", CARD, WHITE),
    ("↓",                            BG,    GRAY),
    ("SQLite 저장 (선택)",            CARD,  WHITE),
    ("↓",                            BG,    GRAY),
    ("JSON 응답",                     CARD,  WHITE),
]

fy = 1.7
for ftext, ffill, fcolor in flow_items:
    fh = 0.28 if ftext == "↓" else 0.48
    add_rect(slide9, 6.5, fy, 6.3, fh, ffill,
             ftext, font_size=11 if ftext != "↓" else 13,
             font_color=fcolor, align=PP_ALIGN.CENTER)
    fy += fh + 0.02


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Database
# ════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide10, BG)
add_title(slide10, "데이터베이스 설계 (SQLite)")

db_tables = [
    (0.3,   3.8, "users",          LIME, [
        "id (PK)", "username", "password_hash", "created_at"]),
    (4.5,   4.5, "workouts",       BLUE, [
        "id (PK)", "user_id (FK)", "exercise_type",
        "exercise_count", "avg_score", "combined_score",
        "dtw_score", "grade", "duration", "created_at"]),
    (9.3,   3.7, "workout_errors", GREEN2, [
        "id (PK)", "workout_id (FK)",
        "error_message", "count", "phase"]),
]

for tx, tw, tname, tcolor, tfields in db_tables:
    add_rect(slide10, tx, 1.3, tw, 4.2, CARD)
    add_rect(slide10, tx, 1.3, tw, 0.42, tcolor)
    add_textbox(slide10, tx + 0.1, 1.34, tw - 0.2, 0.36,
                tname, font_size=14, font_color=BLACK, bold=True,
                align=PP_ALIGN.CENTER)
    add_multiline_textbox(slide10, tx + 0.15, 1.82, tw - 0.3,
                          3.6, tfields, font_size=12, font_color=WHITE)

# Bottom note
add_textbox(slide10, 0.4, 5.75, 12.5, 0.6,
            "• WAL 모드로 동시 접근 지원  • bcrypt 해시 인증  • 유저별 통계 집계 (총 운동, 평균 점수, 선호 운동)",
            font_size=11, font_color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Frontend UI
# ════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide11, BG)
add_title(slide11, "프론트엔드 UI (React + TypeScript)")

pages = [
    (0.3,  "🏠 Home",        [
        "히어로 섹션", "기능 소개 카드", "로그인/로그아웃", "CTA 버튼"]),
    (2.8,  "🏋️ 운동 선택",   [
        "푸시업 / 풀업", "그립 타입 선택",
        "(오버핸드/언더핸드", "/와이드)"]),
    (5.3,  "📤 영상 업로드",  [
        "메인 영상", "레퍼런스 영상",
        "FPS 슬라이더", "분석 로딩 화면"]),
    (7.8,  "📊 결과 대시보드",[
        "등급 / 점수", "프레임별 뷰어",
        "Phase 차트", "AI 피드백"]),
    (10.3, "👤 마이페이지",   [
        "운동 기록", "성장 추이", "통계 요약"]),
]

page_colors = [LIME, BLUE, GREEN2, ORANGE, PURPLE]
for i, (px, ptitle, pbody) in enumerate(pages):
    add_card(slide11, px, 1.4, 2.3, 3.8, ptitle, pbody,
             title_color=page_colors[i], body_font_size=11)

# Tech badges
badges = [
    ("React 18",    BLUE),
    ("TypeScript",  BLUE),
    ("Tailwind CSS",GREEN2),
    ("shadcn/ui",   PURPLE),
    ("React Router",ORANGE),
    ("Vite",        LIME),
]
bx = 0.4
for btext, bcolor in badges:
    bw = 1.95
    add_rect(slide11, bx, 5.5, bw, 0.45, CARD,
             btext, font_size=11, font_color=bcolor, bold=True)
    bx += bw + 0.12


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Gemini AI Feedback
# ════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide12, BG)
add_title(slide12, "Gemini AI 종합 피드백")

# Left: flow
add_textbox(slide12, 0.4, 1.2, 5.5, 0.4,
            "피드백 생성 프로세스", font_size=14, font_color=LIME, bold=True)

flow_boxes = [
    ("분석 결과 수집",           "avg_score, phase별 점수, top 3 오류, DTW 점수", CARD),
    ("구조화된 프롬프트 생성",   "운동 종목, 횟수, 오류 빈도, Phase 약점",         CARD),
    ("Gemini API 호출",          "temperature=0.4, max_tokens=800",               CARD),
    ("마크다운 정리 후 표시",    "",                                               CARD),
]

fy = 1.7
for ftitle, fdetail, ffill in flow_boxes:
    add_rect(slide12, 0.4, fy, 5.6, 0.45, ffill,
             ftitle, font_size=12, font_color=LIME, bold=True,
             align=PP_ALIGN.LEFT)
    if fdetail:
        add_textbox(slide12, 0.55, fy + 0.47, 5.3, 0.3,
                    fdetail, font_size=10, font_color=GRAY)
        add_textbox(slide12, 2.7, fy + 0.9, 0.6, 0.3,
                    "↓", font_size=12, font_color=GRAY, align=PP_ALIGN.CENTER)
        fy += 1.2
    else:
        fy += 0.55

# Right: example feedback
add_textbox(slide12, 6.3, 1.2, 6.7, 0.4,
            "피드백 예시", font_size=14, font_color=BLUE, bold=True)

add_rect(slide12, 6.3, 1.7, 6.6, 4.7, CARD)
example_lines = [
    "총평: 전반적으로 B등급 수준의 자세를 보여주셨습니다.",
    "",
    "주요 개선 사항:",
    "① 팔꿈치 벌림 — 어깨 부상 위험",
    "② 고개 숙임 — 경추 압박 주의",
    "",
    "잘 된 부분:",
    "✓ 등 직선 유지 양호",
    "✓ 손 위치 일정",
    "",
    "권장 드릴: 벽 푸시업으로...",
]
add_multiline_textbox(slide12, 6.45, 1.85, 6.3, 4.4,
                      example_lines, font_size=11.5, font_color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Challenges & Solutions
# ════════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide13, BG)
add_title(slide13, "기술적 도전과 해결 과정")

challenges = [
    (1.3,  "카메라 위치에 따른 DTW 피처 왜곡",
           "body_sway: 절대 x좌표 → 목 대비 상대 좌표"),
    (2.5,  "이진 점수로 미세한 자세 차이 미반영",
           "Soft scoring 도입: 임계값 근방 선형 보간"),
    (3.7,  "정지/준비 구간이 점수에 포함돼 왜곡",
           "ML + Rule-based 하이브리드 활성 구간 필터"),
    (4.9,  "Phase 전환 시 노이즈(jitter) 발생",
           "Hysteresis + 속도 스무딩으로 안정적 감지"),
]

for ry, prob, sol in challenges:
    # Left border accent for problem
    add_rect(slide13, 0.3, ry, 0.07, 0.9, ORANGE)
    add_rect(slide13, 0.37, ry, 5.43, 0.9, CARD)
    add_textbox(slide13, 0.5, ry + 0.15, 5.2, 0.6,
                prob, font_size=12, font_color=WHITE)

    # Arrow
    add_textbox(slide13, 6.0, ry + 0.2, 0.8, 0.5,
                "→", font_size=20, font_color=LIME, bold=True,
                align=PP_ALIGN.CENTER)

    # Right border accent for solution
    add_rect(slide13, 6.9, ry, 0.07, 0.9, LIME)
    add_rect(slide13, 6.97, ry, 5.83, 0.9, CARD)
    add_textbox(slide13, 7.1, ry + 0.15, 5.6, 0.6,
                sol, font_size=12, font_color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Demo Results
# ════════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide14, BG)
add_title(slide14, "분석 결과 화면 구성")

result_cards = [
    (0.4,  1.4, "📈 종합 점수 & 등급", LIME, [
        "Combined Score = avg × 0.7 + DTW × 0.3",
        "S(≥90%) / A(≥70%) / B(≥50%) / C(<50%)"]),
    (6.6,  1.4, "🎞️ 프레임별 뷰어",   BLUE, [
        "Phase 필터링으로 구간별 확인",
        "오류 프레임 하이라이트",
        "스켈레톤 오버레이 표시"]),
    (0.4,  3.9, "📊 Phase 분석 차트",  GREEN2, [
        "Phase별 평균 점수 바차트",
        "약점 Phase 시각적 확인",
        "프레임 분포 히스토그램"]),
    (6.6,  3.9, "🤖 AI 피드백 & PDF", PURPLE, [
        "Gemini AI 종합 코멘트",
        "PDF 리포트 다운로드",
        "운동 기록 자동 저장"]),
]

for rx, ry, rtitle, rcolor, rbody in result_cards:
    add_card(slide14, rx, ry, 5.8, 2.3, rtitle, rbody,
             title_color=rcolor, body_font_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ════════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide15, BG)
add_title(slide15, "마무리")

# Left: learned
add_card(slide15, 0.4, 1.3, 5.8, 4.0,
         "프로젝트에서 배운 것", [
             "• YOLO 포즈 추정 파이프라인 전체 구현",
             "• DTW 알고리즘의 실전 적용",
             "• 통계 기반 시스템 설계 (Cohen's d)",
             "• FastAPI + React 풀스택 개발",
             "• ML 모델과 규칙 기반의 하이브리드 설계",
             "• 스포츠 과학 도메인 지식 습득",
             "• Git 협업 워크플로우",
         ], title_color=LIME, body_font_size=12)

# Right: future
add_card(slide15, 7.0, 1.3, 5.9, 4.0,
         "향후 발전 방향", [
             "• 실시간 웹캠 분석 지원",
             "• 스쿼트, 데드리프트 종목 확장",
             "• 레퍼런스 영상 DB 구축",
             "• 모바일 앱 지원 (React Native)",
             "• 개인화 트레이닝 프로그램 추천",
             "• 클라우드 배포 (AWS / GCP)",
         ], title_color=BLUE, body_font_size=12)

# Bottom banner
add_rect(slide15, 0.4, 5.6, 12.5, 0.8, LIME,
         "PoseCoach — 데이터로 만드는 더 나은 운동  |  28기 프로젝트",
         font_size=18, font_color=BLACK, bold=True)


# ── Save ─────────────────────────────────────────────────────────
OUTPUT = "/Users/sanghooh/newpro/28th-project-posecoach/PoseCoach_발표.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
